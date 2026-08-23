import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# 配置路径
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE_PATH = os.path.join(BASE_DIR, '2026-1-10.pptx')
OUTPUT_PATH = os.path.join(BASE_DIR, 'Judah_Kings_Gen.pptx')

def create_kings_ppt():
    if not os.path.exists(TEMPLATE_PATH):
        print(f"Error: Template not found at {TEMPLATE_PATH}")
        return

    # 1. 加载模板
    print(f"Loading template from {TEMPLATE_PATH}...")
    prs = Presentation(TEMPLATE_PATH)

    # 数据：犹大末期列王
    # 格式: (名字, 好坏, 在位时间, 结局/更迭原因)
    kings_part1 = [
        ("玛拿西 (Manasseh)", "恶 (极恶)", "55年", "希西家之子。重建偶像，引诱百姓犯罪。晚年被掳巴比伦后悔改，归回耶路撒冷善终。"),
        ("亚们 (Amon)", "恶", "2年", "玛拿西之子。效法父亲恶行，未曾自卑。被臣仆在宫中杀害。"),
        ("约西亚 (Josiah)", "好 (极好)", "31年", "亚们之子，8岁登基。得律法书，推行彻底的宗教改革。干预埃及与巴比伦战事，被埃及法老尼哥杀于米吉多。")
    ]

    kings_part2 = [
        ("约哈斯 (Jehoahaz)", "恶", "3个月", "约西亚之子。被埃及法老尼哥废黜，囚禁并带到埃及，死在那里。"),
        ("约雅敬 (Jehoiakim)", "恶", "11年", "约西亚之子，约哈斯之兄。埃及傀儡，后服侍巴比伦。奢华宴乐，压榨百姓。死于任内（或被巴比伦带走途中）。"),
        ("约雅斤 (Jehoiachin)", "恶", "3个月", "约雅敬之子。巴比伦王尼布甲尼撒围城，主动投降。被掳至巴比伦（第一次大被掳，597 BC）。"),
        ("西底家 (Zedekiah)", "恶", "11年", "约西亚之子，约雅斤叔叔。发誓效忠巴比伦后背叛。耶路撒冷被攻破，双眼被剜，被掳至巴比伦（亡国，586 BC）。")
    ]

    # 定义一个辅助函数来添加表格幻灯片
    def add_table_slide(slide_title, data):
        # 尝试使用索引1的版式 (通常是 Title and Content)
        # 如果模板版式很特殊，可能需要调整这个索引
        slide_layout = prs.slide_layouts[1] 
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置标题
        title = slide.shapes.title
        title.text = slide_title

        # 添加表格
        rows = len(data) + 1
        cols = 4
        left = Inches(0.5)
        top = Inches(2.0)
        width = Inches(9.0)
        height = Inches(0.8)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # 设置列宽
        table.columns[0].width = Inches(2.0) # 名字
        table.columns[1].width = Inches(1.2) # 评价
        table.columns[2].width = Inches(1.0) # 时间
        table.columns[3].width = Inches(4.8) # 结局

        # 表头
        headers = ["君王", "评价", "在位", "更迭原因/结局"]
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            # 简单样式设置
            if cell.text_frame.text:
                para = cell.text_frame.paragraphs[0]
                para.font.bold = True
                para.font.size = Pt(14)

        # 填充数据
        for r_idx, row_data in enumerate(data):
            for c_idx, cell_value in enumerate(row_data):
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = cell_value
                
                # 设置字体大小
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(12)
                    if c_idx == 1: # 评价列居中
                        paragraph.alignment = PP_ALIGN.CENTER
    
    # 2. 生成两页幻灯片
    add_table_slide("犹大列王纪 (1)：玛拿西至约西亚", kings_part1)
    add_table_slide("犹大列王纪 (2)：帝国的黄昏", kings_part2)

    # 3. 保存
    print(f"Saving presentation to {OUTPUT_PATH}...")
    prs.save(OUTPUT_PATH)
    print("Done.")

if __name__ == "__main__":
    create_kings_ppt()




