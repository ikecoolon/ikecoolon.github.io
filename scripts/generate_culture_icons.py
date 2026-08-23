import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# 配置路径
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUTPUT_PATH = os.path.join(BASE_DIR, 'Icons_Jew_Greek_Roman.pptx')

# 颜色定义
COLOR_JEWISH_BLUE = RGBColor(0, 56, 184)      # 以色列蓝
COLOR_GREEK_BLUE = RGBColor(13, 94, 175)      # 希腊蓝
COLOR_ROMAN_RED = RGBColor(184, 18, 50)       # 罗马红
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_GOLD = RGBColor(218, 165, 32)

def create_icons_ppt():
    prs = Presentation()
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    p = title.text_frame.paragraphs[0]
    p.text = "示意图标：犹太人、希腊人、罗马人"
    p.font.size = Pt(32)
    p.font.bold = True

    # 辅助函数：绘制圆形背景
    def draw_circle_base(x, y, size, color, label):
        # 外圆
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
        oval.fill.solid()
        oval.fill.fore_color.rgb = color
        oval.line.color.rgb = COLOR_WHITE
        oval.line.width = Pt(3)
        
        # 下方标签
        lbl = slide.shapes.add_textbox(x, y + size + Inches(0.1), size, Inches(0.5))
        p = lbl.text_frame.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        return oval

    y_pos = Inches(2.5)
    size = Inches(2.5)
    gap = Inches(0.5)
    
    # 1. 犹太人 (Jewish) - 代表：六芒星
    x1 = Inches(1.0)
    draw_circle_base(x1, y_pos, size, COLOR_JEWISH_BLUE, "犹太人\n(信仰/律法)")
    
    # 绘制六芒星 (用两个三角形叠加)
    star_size = size * 0.6
    star_x = x1 + (size - star_size) / 2
    star_y = y_pos + (size - star_size) / 2
    
    # 正三角
    t1 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, star_x, star_y, star_size, star_size)
    t1.fill.background() # 无填充 (透明)
    t1.line.color.rgb = COLOR_WHITE
    t1.line.width = Pt(4)
    
    # 倒三角
    t2 = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, star_x, star_y, star_size, star_size)
    t2.rotation = 180
    t2.fill.background()
    t2.line.color.rgb = COLOR_WHITE
    t2.line.width = Pt(4)

    # 2. 希腊人 (Greek) - 代表：Alpha & Omega 或 柱子
    x2 = x1 + size + gap
    draw_circle_base(x2, y_pos, size, COLOR_GREEK_BLUE, "希腊人\n(智慧/文化)")
    
    # 绘制符号 (这里用文本 Α Ω)
    tb_g = slide.shapes.add_textbox(x2, y_pos + size*0.25, size, size*0.5)
    p = tb_g.text_frame.paragraphs[0]
    p.text = "Α Ω" # Alpha Omega
    p.font.size = Pt(60)
    p.font.color.rgb = COLOR_WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # 3. 罗马人 (Roman) - 代表：SPQR (元老院与罗马人民)
    x3 = x2 + size + gap
    draw_circle_base(x3, y_pos, size, COLOR_ROMAN_RED, "罗马人\n(政权/秩序)")
    
    # 绘制符号 (文本 SPQR)
    # 加一个圆环装饰（桂冠意象）
    wreath_size = size * 0.8
    wx = x3 + (size - wreath_size) / 2
    wy = y_pos + (size - wreath_size) / 2
    
    # 内圆圈 (模拟月桂冠)
    wreath = slide.shapes.add_shape(MSO_SHAPE.ARC, wx, wy, wreath_size, wreath_size)
    wreath.line.color.rgb = COLOR_GOLD
    wreath.line.width = Pt(6)
    
    tb_r = slide.shapes.add_textbox(x3, y_pos + size*0.35, size, size*0.3)
    p = tb_r.text_frame.paragraphs[0]
    p.text = "SPQR"
    p.font.name = "Times New Roman"
    p.font.size = Pt(40)
    p.font.color.rgb = COLOR_GOLD
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    print(f"Saving icons to {OUTPUT_PATH}...")
    prs.save(OUTPUT_PATH)
    print("Done.")

if __name__ == "__main__":
    create_icons_ppt()

