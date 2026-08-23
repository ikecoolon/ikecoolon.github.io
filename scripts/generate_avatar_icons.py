import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# 配置路径
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUTPUT_PATH = os.path.join(BASE_DIR, 'Icons_People_Avatar.pptx')

def create_avatar_ppt():
    prs = Presentation()
    slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title.text_frame.text = "人物形象示意图"
    
    y_pos = Inches(2.5)
    size = Inches(2.5) # 圆的直径
    gap = Inches(0.5)
    
    # 基础肤色
    SKIN_COLOR = RGBColor(255, 224, 189) 

    # 辅助：画底圆
    def draw_base(x, color, label):
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y_pos, size, size)
        oval.fill.solid()
        oval.fill.fore_color.rgb = color
        oval.line.fill.background() # 无边框
        
        lbl = slide.shapes.add_textbox(x, y_pos + size + Inches(0.1), size, Inches(0.5))
        p = lbl.text_frame.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(18)
        return x + size/2, y_pos + size/2 # 返回圆心坐标

    # 1. 犹太人 (Jewish)
    # 特征：大胡子，头巾/圆帽
    cx1, cy1 = draw_base(Inches(1.0), RGBColor(230, 240, 255), "犹太人")
    
    # 头
    head_size = Inches(1.0)
    head = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx1 - head_size/2, cy1 - head_size/2, head_size, head_size)
    head.fill.solid()
    head.fill.fore_color.rgb = SKIN_COLOR
    head.line.fill.background()
    
    # 胡子 (倒三角/梯形)
    beard_w = Inches(0.8)
    beard_h = Inches(0.6)
    beard = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, cx1 - beard_w/2, cy1 + Inches(0.1), beard_w, beard_h)
    beard.rotation = 180
    beard.fill.solid()
    beard.fill.fore_color.rgb = RGBColor(80, 50, 20) # 深棕色
    beard.line.fill.background()

    # 帽子 (Kippah - 小半圆)
    hat_w = Inches(0.6)
    hat_h = Inches(0.3)
    hat = slide.shapes.add_shape(MSO_SHAPE.CHORD, cx1 - hat_w/2, cy1 - head_size/2 - Inches(0.05), hat_w, hat_h)
    hat.rotation = 180 # 扣在头上
    hat.fill.solid()
    hat.fill.fore_color.rgb = RGBColor(0, 0, 100) # 深蓝
    hat.line.fill.background()


    # 2. 希腊人 (Greek)
    # 特征：卷发/桂冠，白袍
    cx2, cy2 = draw_base(Inches(1.0) + size + gap, RGBColor(240, 255, 240), "希腊人")
    
    # 身体 (长袍 - 梯形)
    body_w = Inches(1.4)
    body_h = Inches(1.2)
    body = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, cx2 - body_w/2, cy2 + Inches(0.3), body_w, body_h)
    body.fill.solid()
    body.fill.fore_color.rgb = RGBColor(255, 255, 255) # 白袍
    body.line.color.rgb = RGBColor(200, 200, 200)

    # 头
    head2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx2 - head_size/2, cy2 - head_size/2, head_size, head_size)
    head2.fill.solid()
    head2.fill.fore_color.rgb = SKIN_COLOR
    head2.line.fill.background()

    # 桂冠 (绿色圆环片段 或 两个小叶子)
    leaf_w = Inches(0.3)
    leaf_h = Inches(0.15)
    # 左叶
    leaf1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx2 - head_size/2, cy2 - head_size/3, leaf_w, leaf_h)
    leaf1.rotation = 30
    leaf1.fill.solid()
    leaf1.fill.fore_color.rgb = RGBColor(34, 139, 34) # 森林绿
    leaf1.line.fill.background()
    # 右叶
    leaf2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx2 + head_size/2 - leaf_w, cy2 - head_size/3, leaf_w, leaf_h)
    leaf2.rotation = -30
    leaf2.fill.solid()
    leaf2.fill.fore_color.rgb = RGBColor(34, 139, 34)
    leaf2.line.fill.background()


    # 3. 罗马人 (Roman)
    # 特征：头盔 (Helmet)，红缨
    cx3, cy3 = draw_base(Inches(1.0) + (size + gap)*2, RGBColor(255, 240, 240), "罗马人")

    # 身体 (铠甲 - 红色/金色)
    body3 = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, cx3 - body_w/2, cy3 + Inches(0.3), body_w, body_h)
    body3.fill.solid()
    body3.fill.fore_color.rgb = RGBColor(178, 34, 34) # 罗马红
    body3.line.fill.background()

    # 头
    head3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx3 - head_size/2, cy3 - head_size/2, head_size, head_size)
    head3.fill.solid()
    head3.fill.fore_color.rgb = SKIN_COLOR
    head3.line.fill.background()

    # 头盔主体 (灰色半圆)
    helmet_w = Inches(1.1)
    helmet_h = Inches(0.6)
    helmet = slide.shapes.add_shape(MSO_SHAPE.CHORD, cx3 - helmet_w/2, cy3 - head_size/2, helmet_w, helmet_h)
    helmet.rotation = 180
    helmet.fill.solid()
    helmet.fill.fore_color.rgb = RGBColor(169, 169, 169) # 银灰
    helmet.line.fill.background()

    # 头盔红缨 (扫帚状 - 扇形或梯形)
    crest_w = Inches(0.8)
    crest_h = Inches(0.4)
    crest = slide.shapes.add_shape(MSO_SHAPE.ARC, cx3 - crest_w/2, cy3 - head_size - Inches(0.1), crest_w, crest_h)
    # ARC 只有线条，换一个实心形状
    crest = slide.shapes.add_shape(MSO_SHAPE.MOON, cx3 - crest_w/2, cy3 - head_size/1.5 - Inches(0.2), crest_w, crest_h)
    crest.rotation = 90
    crest.fill.solid()
    crest.fill.fore_color.rgb = RGBColor(220, 20, 60) # 猩红
    crest.line.fill.background()

    print(f"Saving avatars to {OUTPUT_PATH}...")
    prs.save(OUTPUT_PATH)
    print("Done.")

if __name__ == "__main__":
    create_avatar_ppt()

