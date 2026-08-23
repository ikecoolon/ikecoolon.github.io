import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt

# 配置路径
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
TEMPLATE_PATH = os.path.join(BASE_DIR, '2026-1-10.pptx')
OUTPUT_PATH = os.path.join(BASE_DIR, 'Judah_Kings_Timeline_Josiah_OneSlide_Callouts.pptx')


def _get_blank_layout(prs: Presentation):
    """
    Pick a blank slide layout if possible (to maximize usable area),
    falling back gracefully if the template doesn't have one.
    """
    for layout in prs.slide_layouts:
        try:
            if len(layout.placeholders) == 0:
                return layout
        except Exception:
            continue
    return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]


def _set_shape_fill_theme(shape, theme_color, *, brightness=None):
    shape.fill.solid()
    shape.fill.fore_color.theme_color = theme_color
    if brightness is not None:
        shape.fill.fore_color.brightness = brightness


def _set_shape_line_theme(shape, theme_color, *, width_pt=2, brightness=None):
    shape.line.width = Pt(width_pt)
    shape.line.color.theme_color = theme_color
    if brightness is not None:
        shape.line.color.brightness = brightness


def _set_font_theme_color(font, theme_color, *, brightness=None):
    font.color.theme_color = theme_color
    if brightness is not None:
        font.color.brightness = brightness

def create_timeline_ppt():
    if not os.path.exists(TEMPLATE_PATH):
        print(f"Error: Template not found at {TEMPLATE_PATH}")
        return

    print(f"Loading template from {TEMPLATE_PATH}...")
    prs = Presentation(TEMPLATE_PATH)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # 单页：从约西亚开始
    kings = [
        {"name": "约西亚", "eval": "好王", "years": 31.0, "time": "31年"},
        {"name": "约哈斯", "eval": "坏王", "years": 0.25, "time": "3个月"},
        {"name": "约雅敬", "eval": "坏王", "years": 11.0, "time": "11年"},
        {"name": "约雅斤", "eval": "坏王", "years": 0.25, "time": "3个月"},
        {"name": "西底家", "eval": "坏王", "years": 11.0, "time": "11年"},
    ]

    # 进程图需要“按时长拉伸”，但 3 个月的王必须保留最小可读宽度（等价年数）
    MIN_YEARS_EQUIV = 2.0
    weights = [max(k["years"], MIN_YEARS_EQUIV) for k in kings]
    total_weight = sum(weights)

    slide = prs.slides.add_slide(_get_blank_layout(prs))

    margin_x = Inches(0.3)
    margin_top = Inches(0.15)
    margin_bottom = Inches(0.15)
    gap = Inches(0.12)
    callout_box_h = Inches(0.85)
    callout_gap = Inches(0.10)

    # 标题（留白最少，尽量让图撑满）
    title_h = Inches(0.75)
    title = slide.shapes.add_textbox(margin_x, margin_top, slide_w - 2 * margin_x, title_h)
    tf = title.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "犹大末期：约西亚 → 亡国（并但以理/三次被掳）"
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(28)
    _set_font_theme_color(p.font, MSO_THEME_COLOR.TEXT_1)

    p2 = tf.add_paragraph()
    p2.text = "绿=好王   橙=坏王"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(16)
    p2.font.bold = True
    _set_font_theme_color(p2.font, MSO_THEME_COLOR.TEXT_1)

    # 时间条（按在位时长伸缩）
    # 在标题与时间条之间预留一行“更迭原因”标注（上方）
    bar_y = margin_top + title_h + gap + callout_box_h + callout_gap
    bar_h = Inches(1.55)  # 矮一点，给上下标注让空间
    timeline_x = margin_x
    timeline_w = slide_w - 2 * margin_x

    segments = []
    x = timeline_x
    for idx, (k, w) in enumerate(zip(kings, weights)):
        if idx == len(kings) - 1:
            seg_w = (timeline_x + timeline_w) - x
        else:
            seg_w = int(timeline_w * (w / total_weight))
        segments.append({"king": k, "x": x, "w": seg_w})
        x += seg_w

    for seg in segments:
        k = seg["king"]
        seg_x = seg["x"]
        seg_w = seg["w"]
        is_good = k["eval"].startswith("好")
        # 模板主题配色：accent6=绿，accent2=橙，accent1=蓝，accent3=灰
        theme = MSO_THEME_COLOR.ACCENT_6 if is_good else MSO_THEME_COLOR.ACCENT_2

        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, seg_x, bar_y, seg_w, bar_h)
        _set_shape_fill_theme(rect, theme, brightness=0.75)
        _set_shape_line_theme(rect, theme, width_pt=3)

        # 段内只放“王名”（避免短在位导致文字挤压）。太窄的段用旋转文字保证可读。
        narrow = seg_w < Inches(0.75)
        if not narrow:
            tb = slide.shapes.add_textbox(seg_x, bar_y, seg_w, bar_h)
            tf = tb.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            p1 = tf.paragraphs[0]
            p1.text = k["name"]
            p1.alignment = PP_ALIGN.CENTER
            p1.font.bold = True
            p1.font.size = Pt(40)
            _set_font_theme_color(p1.font, MSO_THEME_COLOR.TEXT_1)
        else:
            cx = seg_x + (seg_w // 2)
            cy = bar_y + (bar_h // 2)
            tb_w = bar_h
            tb_h = seg_w
            tb = slide.shapes.add_textbox(cx - (tb_w // 2), cy - (tb_h // 2), tb_w, tb_h)
            tb.rotation = 90
            tf = tb.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE

            p1 = tf.paragraphs[0]
            p1.text = k["name"]
            p1.alignment = PP_ALIGN.CENTER
            p1.font.bold = True
            p1.font.size = Pt(24)
            _set_font_theme_color(p1.font, MSO_THEME_COLOR.TEXT_1)

    # 更迭原因：用“文本框 + 箭头”贴近时间条的上/下方，指向交界处
    def add_reason_callout(*, anchor_x, is_above, text, box_x):
        box_w = Inches(3.1)
        box_h = callout_box_h
        # 这里的箭头头部要能塞进 callout_gap（否则会压到文本框/时间条）
        tri_w = Inches(0.20)
        tri_h = Inches(0.08)

        # 二次兜底：确保 box_x 在画布范围内
        min_x = timeline_x
        max_x = timeline_x + timeline_w - box_w
        if box_x < min_x:
            box_x = min_x
        if box_x > max_x:
            box_x = max_x
        box_center_x = box_x + (box_w // 2)

        if is_above:
            box_y = margin_top + title_h + gap
            tri_rot = 180  # 指向下
            tri_y = bar_y - tri_h
            # 箭头：从文本框底部连到时间条（具体王所在段的顶部）
            line_x1, line_y1 = box_center_x, box_y + box_h
            line_x2, line_y2 = anchor_x, bar_y
        else:
            box_y = bar_y + bar_h + callout_gap
            tri_rot = 0  # 指向上
            tri_y = bar_y + bar_h
            # 箭头：从文本框顶部连到时间条（具体王所在段的底部）
            line_x1, line_y1 = box_center_x, box_y
            line_x2, line_y2 = anchor_x, bar_y + bar_h

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, box_x, box_y, box_w, box_h)
        # 浅色底 + 灰边，确保投影屏幕可读
        _set_shape_fill_theme(box, MSO_THEME_COLOR.BACKGROUND_1)
        _set_shape_line_theme(box, MSO_THEME_COLOR.ACCENT_3, width_pt=2)

        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(18)
        _set_font_theme_color(p.font, MSO_THEME_COLOR.TEXT_1)

        # 连接线（python-pptx 1.0.2 无箭头端点：用三角形模拟箭头）
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, line_x1, line_y1, line_x2, line_y2)
        conn.line.width = Pt(2)
        conn.line.color.theme_color = MSO_THEME_COLOR.ACCENT_3

        tri = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            anchor_x - (tri_w // 2),
            tri_y,
            tri_w,
            tri_h,
        )
        tri.rotation = tri_rot
        _set_shape_fill_theme(tri, MSO_THEME_COLOR.ACCENT_3)
        tri.line.fill.background()

    def _pack_callouts(items, box_w):
        """Prevent overlaps in the same row by packing left-to-right, then correcting from the right edge."""
        if not items:
            return
        row_gap = Inches(0.10)
        min_x = timeline_x
        max_x = timeline_x + timeline_w - box_w

        items.sort(key=lambda d: d["anchor_x"])

        # initial placement near anchor
        for it in items:
            bx = it["anchor_x"] - (box_w // 2)
            if bx < min_x:
                bx = min_x
            if bx > max_x:
                bx = max_x
            it["box_x"] = bx

        # forward sweep
        for i in range(1, len(items)):
            prev_end = items[i - 1]["box_x"] + box_w
            if items[i]["box_x"] < prev_end + row_gap:
                items[i]["box_x"] = prev_end + row_gap

        # clamp to right edge + backward sweep
        if items[-1]["box_x"] > max_x:
            items[-1]["box_x"] = max_x
            for i in range(len(items) - 2, -1, -1):
                allowed = items[i + 1]["box_x"] - row_gap - box_w
                if items[i]["box_x"] > allowed:
                    items[i]["box_x"] = allowed

            # ensure left edge
            if items[0]["box_x"] < min_x:
                shift = min_x - items[0]["box_x"]
                for it in items:
                    it["box_x"] += shift
                for i in range(1, len(items)):
                    prev_end = items[i - 1]["box_x"] + box_w
                    if items[i]["box_x"] < prev_end + row_gap:
                        items[i]["box_x"] = prev_end + row_gap

    def _seg_center(seg):
        return seg["x"] + (seg["w"] // 2)

    # 更迭原因：绑定到“具体王”的长方形段（箭头连到该王）
    callouts = [
        {"seg_idx": 1, "text": "约哈斯继位：约西亚战死（米吉多）", "is_above": True},
        {"seg_idx": 2, "text": "约雅敬继位：埃及废/掳约哈斯", "is_above": False},
        {"seg_idx": 3, "text": "约雅斤继位：约雅敬死", "is_above": True},
        {"seg_idx": 4, "text": "西底家继位：597 投降被掳", "is_above": False},
        {"seg_idx": 4, "text": "西底家末：叛巴比伦 → 586 亡国", "is_above": True},
    ]

    above_items = []
    below_items = []
    for c in callouts:
        seg = segments[c["seg_idx"]]
        it = {"anchor_x": _seg_center(seg), "text": c["text"], "is_above": c["is_above"]}
        (above_items if c["is_above"] else below_items).append(it)

    box_w = Inches(3.1)
    _pack_callouts(above_items, box_w)
    _pack_callouts(below_items, box_w)

    for it in above_items + below_items:
        add_reason_callout(anchor_x=it["anchor_x"], is_above=it["is_above"], text=it["text"], box_x=it["box_x"])

    # 但以理：从 605（约雅敬第三年左右）开始的横向条
    daniel_h = Inches(0.75)
    # 下方还要留出“更迭原因”标注一行
    daniel_y = bar_y + bar_h + callout_gap + callout_box_h + gap
    jego = segments[2]  # 约雅敬
    daniel_start = int(jego["x"] + jego["w"] * 0.36)  # 605 大约处于约雅敬在位早期
    daniel_end = timeline_x + timeline_w
    daniel_w = daniel_end - daniel_start

    band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, daniel_start, daniel_y, daniel_w, daniel_h)
    _set_shape_fill_theme(band, MSO_THEME_COLOR.ACCENT_1, brightness=0.85)
    _set_shape_line_theme(band, MSO_THEME_COLOR.ACCENT_1, width_pt=3)

    tf = band.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "但以理：605 被掳 → 跨越 3 王（约雅敬 → 约雅斤 → 西底家）"
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(20)
    _set_font_theme_color(p.font, MSO_THEME_COLOR.ACCENT_1)

    # 底部：仅保留“被掳节点”，字体更大，方便现场观众看清
    bottom_y = daniel_y + daniel_h + gap
    bottom_h = slide_h - bottom_y - margin_bottom

    ex = slide.shapes.add_textbox(timeline_x, bottom_y, timeline_w, bottom_h)
    tf = ex.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = "巴比伦掳掠（南国）"
    p.font.bold = True
    p.font.size = Pt(22)
    _set_font_theme_color(p.font, MSO_THEME_COLOR.ACCENT_1)
    p.alignment = PP_ALIGN.CENTER

    for t in [
        "605 第一次被掳：但以理",
        "597 第二次被掳：王族/精英（以西结同被掳）",
        "586 第三次被掳：亡国 / 圣殿被毁",
    ]:
        pp = tf.add_paragraph()
        pp.text = t
        pp.font.size = Pt(20)
        pp.font.bold = True
        _set_font_theme_color(pp.font, MSO_THEME_COLOR.TEXT_1)
        pp.space_before = Pt(10)
        pp.alignment = PP_ALIGN.CENTER

    print(f"Saving presentation to {OUTPUT_PATH}...")
    prs.save(OUTPUT_PATH)
    print("Done.")

if __name__ == "__main__":
    create_timeline_ppt()
