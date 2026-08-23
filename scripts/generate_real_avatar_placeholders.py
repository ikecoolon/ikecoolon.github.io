import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# 配置路径
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUTPUT_PATH = os.path.join(BASE_DIR, 'Icons_Real_Avatar_Placeholder.pptx')

def create_real_avatar_ppt():
    prs = Presentation()
    slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title.text_frame.text = "人物头像 (请右键'更改图片'替换为真人素材)"
    
    y_pos = Inches(2.5)
    size = Inches(2.5) 
    gap = Inches(0.5)
    
    def draw_avatar_placeholder(x, label, emoji_char, color):
        # 1. 外框 (圆形)
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y_pos, size, size)
        oval.fill.solid()
        oval.fill.fore_color.rgb = color # 背景色
        oval.line.width = Pt(0) # 无边框
        
        # 2. 用大号 Emoji 模拟真人头像 (作为临时占位)
        # 注意：PPT 文本框中的 Emoji 渲染取决于系统，通常是彩色的
        tb = slide.shapes.add_textbox(x, y_pos + Inches(0.2), size, size)
        p = tb.text_frame.paragraphs[0]
        p.text = emoji_char
        p.font.size = Pt(130) # 极大字号
        p.alignment = PP_ALIGN.CENTER
        
        # 3. 标签
        lbl = slide.shapes.add_textbox(x, y_pos + size + Inches(0.1), size, Inches(0.8))
        tf = lbl.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(18)
        
        # 4. 提示文字
        p2 = tf.add_paragraph()
        p2.text = "(右键圆圈->更改图片)"
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(100, 100, 100)
        p2.alignment = PP_ALIGN.CENTER

    # 1. 犹太人
    # Emoji: 🧔🏻‍♂️ (浅肤色大胡子) 或 🕍 (但你要人物，用人)
    # 使用通用的 "Person: Beard"
    draw_avatar_placeholder(Inches(1.0), "犹太人", "🧔🏻‍♂️", RGBColor(230, 240, 255))

    # 2. 希腊人
    # Emoji: 👱🏼‍♂️ (金发男) 或 🏛️ (神殿) -> 用 🧐 (单片眼镜/学者感) 或 👱🏼‍♂️
    draw_avatar_placeholder(Inches(1.0) + size + gap, "希腊人", "👱🏼‍♂️", RGBColor(240, 255, 240))

    # 3. 罗马人
    # Emoji: ⛑️ (救援头盔-形似) 或 💂🏼‍♂️ (卫兵) -> 这里没有完美的罗马士兵Emoji
    # 我们可以用 💂🏼‍♂️ 暂代，或者用 ⚔️ 
    # 既然你要真人感，这里放一个带头盔的或者一般的严肃男
    draw_avatar_placeholder(Inches(1.0) + (size + gap)*2, "罗马人", "💂🏼‍♂️", RGBColor(255, 240, 240))

    print(f"Saving avatars to {OUTPUT_PATH}...")
    prs.save(OUTPUT_PATH)
    print("Done.")

if __name__ == "__main__":
    create_real_avatar_ppt()

